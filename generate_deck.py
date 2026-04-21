#!/usr/bin/env python3
"""Generate visual slide deck images from a markdown deck file using Gemini."""

import argparse
import concurrent.futures
import json
import mimetypes
import os
import subprocess
import sys
import threading
import time
from datetime import datetime

from google import genai
from google.genai import types
from pptx import Presentation
from pptx.util import Emu


PARTITION_MODEL = "gemini-2.5-flash-lite"
IMAGE_MODEL = "gemini-3.1-flash-image-preview"

MAX_RETRIES = 3
RETRY_BACKOFF = [5, 15, 30]  # seconds to wait before each retry

# Pricing per million tokens (USD) — update as needed
PRICING = {
    PARTITION_MODEL: {"input": 0.075, "output": 0.30, "thinking": 0.30},
    IMAGE_MODEL: {"input": 0.10, "output": 0.40, "thinking": 0.40, "image_output": 3.90},
}

# ── Terminal styling ─────────────────────────────────────────────────────────

BOLD = "\033[1m"
DIM = "\033[2m"
RED = "\033[31m"
GREEN = "\033[32m"
YELLOW = "\033[33m"
BLUE = "\033[34m"
MAGENTA = "\033[35m"
CYAN = "\033[36m"
WHITE = "\033[37m"
RESET = "\033[0m"

BAR_FILL = "━"
BAR_EMPTY = "╌"
BAR_WIDTH = 30

SPINNER_FRAMES = ["⠋", "⠙", "⠹", "⠸", "⠼", "⠴", "⠦", "⠧", "⠇", "⠏"]


def progress_bar(fraction, width=BAR_WIDTH, color=CYAN):
    """Render a colored progress bar string."""
    fraction = max(0.0, min(1.0, fraction))
    filled = int(width * fraction)
    empty = width - filled
    pct = fraction * 100
    return f"{color}{BAR_FILL * filled}{DIM}{BAR_EMPTY * empty}{RESET} {pct:5.1f}%"


def format_size(nbytes):
    """Human-readable file size."""
    if nbytes >= 1024 * 1024:
        return f"{nbytes / (1024 * 1024):.1f} MB"
    elif nbytes >= 1024:
        return f"{nbytes / 1024:.0f} KB"
    return f"{nbytes} B"


def format_duration(seconds):
    """Human-readable duration."""
    if seconds < 60:
        return f"{seconds:.1f}s"
    m, s = divmod(int(seconds), 60)
    return f"{m}m{s:02d}s"


# ── Logging ──────────────────────────────────────────────────────────────────

_print_lock = threading.Lock()


def log(msg, prefix="INFO", color=WHITE):
    """Thread-safe timestamped log."""
    ts = time.strftime("%H:%M:%S")
    with _print_lock:
        print(f"{DIM}[{ts}]{RESET} {color}{BOLD}[{prefix}]{RESET} {msg}", flush=True)


def log_header(title):
    """Print a prominent section header."""
    line = "═" * 60
    with _print_lock:
        print(f"\n{BOLD}{MAGENTA}{line}{RESET}", flush=True)
        print(f"{BOLD}{MAGENTA}  {title}{RESET}", flush=True)
        print(f"{BOLD}{MAGENTA}{line}{RESET}\n", flush=True)


def log_subheader(title):
    """Print a smaller section divider."""
    with _print_lock:
        print(f"\n{CYAN}{BOLD}── {title} ──{RESET}\n", flush=True)


# ── Cost computation ─────────────────────────────────────────────────────────

def compute_cost(model, usage):
    """Compute USD cost from usage metadata and model pricing."""
    if usage is None:
        return 0.0, {}
    rates = PRICING.get(model, {})

    prompt_tokens = getattr(usage, "prompt_token_count", 0) or 0
    candidates_tokens = getattr(usage, "candidates_token_count", 0) or 0
    thoughts_tokens = getattr(usage, "thoughts_token_count", 0) or 0

    image_output_tokens = 0
    text_output_tokens = candidates_tokens
    details = getattr(usage, "candidates_tokens_details", None)
    if details:
        for d in details:
            if d.modality and d.modality.name == "IMAGE":
                image_output_tokens = d.token_count or 0
            elif d.modality and d.modality.name == "TEXT":
                text_output_tokens = d.token_count or 0

    input_cost = prompt_tokens / 1_000_000 * rates.get("input", 0)
    output_cost = text_output_tokens / 1_000_000 * rates.get("output", 0)
    thinking_cost = thoughts_tokens / 1_000_000 * rates.get("thinking", 0)
    image_cost = image_output_tokens / 1_000_000 * rates.get("image_output", 0)

    total = input_cost + output_cost + thinking_cost + image_cost
    breakdown = {
        "input_tokens": prompt_tokens,
        "output_tokens": text_output_tokens,
        "thinking_tokens": thoughts_tokens,
        "image_output_tokens": image_output_tokens,
        "input_cost": input_cost,
        "output_cost": output_cost,
        "thinking_cost": thinking_cost,
        "image_cost": image_cost,
        "total": total,
    }
    return total, breakdown


# ── Partition prompt ─────────────────────────────────────────────────────────

PARTITION_PROMPT = """\
You are a slide deck parser. Given a document containing multiple pages/slides, \
extract individual text-to-image generation prompts — one per page.

The input may use any of these formats (or a mix):

FORMAT A — Markdown with --- separators:
- A global style directive appears before the first --- (often in a blockquote like > *[SLIDE FORMAT: ...]*)
- Each slide is separated by --- and contains a VISUAL blockquote (> *[VISUAL — ...]*)
- Combine the global style directive with each slide's VISUAL content into one prompt.

FORMAT B — XML-style <page> tags:
- Global style/design-system description appears before the first <page> tag.
- Each page is wrapped in <page number="N" title="..."> ... </page>.
- The full content of each page IS the image generation prompt. Combine it with the global \
style description into one complete prompt.

FORMAT C — Any other structure:
- Look for clear page/slide boundaries (numbered sections, headings, horizontal rules, etc.).
- Identify any shared style or design-system description and prepend it to each page's prompt.

For each page/slide, produce a complete, self-contained image generation prompt that includes \
both the global style parameters AND the full page-specific content. Clean up any markup \
(blockquote markers, XML tags, etc.) so the result is plain text ready for an image model.

If a page has both a detailed description and a condensed version, use the DETAILED description \
as the prompt (it contains more visual information for the image model).

Return a JSON array of objects with these fields:
- slide_number (int): the page/slide number as given in the document (or 1-indexed if unnumbered)
- title (string): the page/slide title or heading
- prompt (string): the complete image generation prompt
- speaker_notes (string or null): any speaking notes found on the page, cleaned up. Null if none.

Return ONLY valid JSON, no markdown fences or extra text.

Here is the document:

"""


# ── Args ─────────────────────────────────────────────────────────────────────

def parse_args():
    parser = argparse.ArgumentParser(
        description="Generate slide deck images from a markdown file using Gemini.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""\
examples:
  %(prog)s deck.md --all                        Generate all slides at 512px
  %(prog)s deck.md --all -r 2k -o ./slides      Generate all at 2K into ./slides
  %(prog)s deck.md --slides 3,5,6,10            Regenerate only slides 3, 5, 6, 10
  %(prog)s deck.md --slides 3+                  Generate slide 3 and all after it
  %(prog)s deck.md --slides 1,2 -r 4k -w 10    Slides 1 & 2 at 4K, 10 workers
  %(prog)s deck.md --all -p 3                   Generate 3 variations per slide
  %(prog)s prompt.md --onesheet                  Generate a single image from prompt
  %(prog)s prompt.md --onesheet -p 5 -r 2k      5 variations of a single prompt at 2K

environment:
  GEMINI_API_KEY    Required. Your Google Gemini API key.
                    Get one at https://aistudio.google.com/apikey
""",
    )
    parser.add_argument("deck_file", help="Path to the markdown deck file")
    parser.add_argument(
        "-o", "--output", default="./output", help="Output directory (default: ./output)"
    )
    parser.add_argument(
        "-r",
        "--resolution",
        choices=["512", "1k", "2k", "4k"],
        default="512",
        help="Image resolution (default: 512)",
    )
    parser.add_argument(
        "-w",
        "--max-parallel",
        type=int,
        default=20,
        help="Max concurrent image generations (default: 20)",
    )
    parser.add_argument(
        "-p",
        "--permutations",
        type=int,
        default=1,
        help="Number of variations to generate per slide (default: 1)",
    )

    # Slide selection: must specify --all, --slides, or --onesheet
    slide_group = parser.add_mutually_exclusive_group(required=True)
    slide_group.add_argument(
        "--all",
        action="store_true",
        help="Generate all slides",
    )
    slide_group.add_argument(
        "--slides",
        type=str,
        help="Comma-separated slide numbers (e.g. 3,5,6,10) or N+ for slide N onward (e.g. 3+)",
    )
    slide_group.add_argument(
        "--onesheet",
        action="store_true",
        help="Treat the input file as a single image prompt (no deck parsing)",
    )

    return parser.parse_args()


def parse_slide_numbers(slides_str, total_slides=None):
    """Parse '3,5,6,10' or '3+' into a set of ints."""
    range_match = slides_str.strip()
    if range_match.endswith("+"):
        start_str = range_match[:-1].strip()
        try:
            start = int(start_str)
        except ValueError:
            print(f"Error: invalid slide range '{slides_str}'", file=sys.stderr)
            sys.exit(1)
        if total_slides is not None:
            return set(range(start, total_slides + 1))
        return {("range_start", start)}

    nums = set()
    for part in slides_str.split(","):
        part = part.strip()
        if not part:
            continue
        try:
            nums.add(int(part))
        except ValueError:
            print(f"Error: invalid slide number '{part}'", file=sys.stderr)
            sys.exit(1)
    return nums


def resolve_slide_range(selected_slides, max_slide_number):
    """Expand range sentinels from parse_slide_numbers into real slide numbers."""
    resolved = set()
    for item in selected_slides:
        if isinstance(item, tuple) and item[0] == "range_start":
            start = item[1]
            resolved.update(range(start, max_slide_number + 1))
        else:
            resolved.add(item)
    return resolved


# ── Deck parsing ─────────────────────────────────────────────────────────────

def partition_deck(client, markdown_content):
    """Use Gemini Flash Lite to parse the deck into individual slide prompts."""
    prompt_text = PARTITION_PROMPT + markdown_content

    log_subheader("Parsing deck")
    log(f"Sending {len(prompt_text):,} chars to {BOLD}{PARTITION_MODEL}{RESET}", "PARSE", BLUE)
    log(f"Prompt: {len(markdown_content):,} chars of deck content + {len(PARTITION_PROMPT):,} chars of instructions", "PARSE", BLUE)

    t0 = time.time()
    chunks_received = 0
    json_text = ""
    usage = None
    last_log_time = t0

    for chunk in client.models.generate_content_stream(
        model=PARTITION_MODEL,
        contents=[
            types.Content(
                role="user",
                parts=[types.Part.from_text(text=prompt_text)],
            ),
        ],
        config=types.GenerateContentConfig(
            response_mime_type="application/json",
            temperature=0.1,
            max_output_tokens=65536,
        ),
    ):
        chunks_received += 1
        if chunk.text:
            json_text += chunk.text
            now = time.time()
            # Log every 2 seconds to avoid spam
            if now - last_log_time >= 2.0 or chunks_received <= 2:
                elapsed = now - t0
                chars_per_sec = len(json_text) / elapsed if elapsed > 0 else 0
                bar = progress_bar(min(len(json_text) / 15000, 1.0), color=BLUE)
                log(
                    f"{bar}  {len(json_text):,} chars received  "
                    f"{DIM}({chars_per_sec:.0f} chars/s, {format_duration(elapsed)}){RESET}",
                    "PARSE", BLUE,
                )
                last_log_time = now
        if getattr(chunk, "usage_metadata", None):
            usage = chunk.usage_metadata

    elapsed = time.time() - t0
    parse_cost, parse_breakdown = compute_cost(PARTITION_MODEL, usage)

    log(f"{GREEN}✓ Parse complete{RESET}  {len(json_text):,} chars in {format_duration(elapsed)}", "PARSE", GREEN)
    if parse_breakdown:
        log(
            f"  Tokens: {DIM}in:{parse_breakdown['input_tokens']:,}  "
            f"out:{parse_breakdown['output_tokens']:,}  "
            f"think:{parse_breakdown['thinking_tokens']:,}{RESET}  "
            f"Cost: {BOLD}${parse_cost:.4f}{RESET}",
            "PARSE", GREEN,
        )

    slides = json.loads(json_text)
    for slide in slides:
        if not all(k in slide for k in ("slide_number", "title", "prompt")):
            raise ValueError(f"Invalid slide entry missing required fields: {slide}")

    log("", "PARSE", WHITE)
    for s in slides:
        log(
            f"  {BOLD}Slide {s['slide_number']:2d}{RESET}  {s['title']}  "
            f"{DIM}({len(s['prompt']):,} char prompt){RESET}",
            "PARSE", WHITE,
        )

    return slides, parse_cost


# ── File helpers ─────────────────────────────────────────────────────────────

_file_lock = threading.Lock()


def next_available_path(output_dir, base_name, ext):
    """Return a non-colliding path, appending -variation-NN if the base name is taken."""
    with _file_lock:
        path = os.path.join(output_dir, f"{base_name}{ext}")
        if not os.path.exists(path):
            open(path, "wb").close()
            return path
        n = 1
        while True:
            path = os.path.join(output_dir, f"{base_name}-variation-{n:02d}{ext}")
            if not os.path.exists(path):
                open(path, "wb").close()
                return path
            n += 1


def save_binary_file(file_name, data):
    with open(file_name, "wb") as f:
        f.write(data)


# ── Slide image generation ───────────────────────────────────────────────────

def is_retryable(exc):
    """Check if an exception is a transient error worth retrying."""
    exc_str = str(exc)
    # HTTP status codes that are transient
    for code in ("429", "500", "502", "503", "504"):
        if code in exc_str:
            return True
    # Common transient error messages
    for msg in ("UNAVAILABLE", "Deadline expired", "deadline", "RESOURCE_EXHAUSTED",
                "overloaded", "rate limit", "timeout", "Timeout", "connection"):
        if msg.lower() in exc_str.lower():
            return True
    return False


def generate_slide_image(client, slide, resolution, output_dir):
    """Generate a single slide image with automatic retry on transient errors."""
    slide_num = slide["slide_number"]
    slide_title = slide["title"]
    tag = f"SLIDE {slide_num:02d}"

    image_size = {"512": "512", "1k": "1K", "2k": "2K", "4k": "4K"}[resolution]
    prompt = slide["prompt"]

    for attempt in range(1, MAX_RETRIES + 1):
        try:
            return _generate_slide_image_once(
                client, slide_num, slide_title, prompt, image_size, output_dir, tag, attempt
            )
        except Exception as e:
            if attempt < MAX_RETRIES and is_retryable(e):
                wait = RETRY_BACKOFF[attempt - 1]
                log(
                    f"{YELLOW}⚠ Attempt {attempt}/{MAX_RETRIES} failed: {e}{RESET}",
                    tag, YELLOW,
                )
                log(
                    f"{YELLOW}  Retrying in {wait}s...{RESET}",
                    tag, YELLOW,
                )
                # Show countdown
                for remaining in range(wait, 0, -1):
                    spinner = SPINNER_FRAMES[remaining % len(SPINNER_FRAMES)]
                    log(f"  {DIM}{spinner} Waiting... {remaining}s{RESET}", tag, YELLOW)
                    time.sleep(1)
                log(f"{CYAN}↻ Retry {attempt + 1}/{MAX_RETRIES} starting{RESET}", tag, CYAN)
            else:
                raise


def _generate_slide_image_once(client, slide_num, slide_title, prompt, image_size, output_dir, tag, attempt):
    """Single attempt at generating a slide image."""
    attempt_label = f" {DIM}(attempt {attempt}/{MAX_RETRIES}){RESET}" if attempt > 1 else ""

    # ── Upload phase ──
    log(f"{CYAN}⬆ Uploading prompt{RESET}  {len(prompt):,} chars  res={image_size}{attempt_label}", tag, CYAN)

    prompt_bar = progress_bar(0.0, color=CYAN)
    log(f"  {prompt_bar}  Sending to {IMAGE_MODEL}...", tag, CYAN)

    contents = [
        types.Content(
            role="user",
            parts=[types.Part.from_text(text=prompt)],
        ),
    ]
    config = types.GenerateContentConfig(
        thinking_config=types.ThinkingConfig(thinking_level="HIGH"),
        response_modalities=["IMAGE", "TEXT"],
        image_config=types.ImageConfig(image_size=image_size),
    )

    prompt_bar = progress_bar(1.0, color=CYAN)
    log(f"  {prompt_bar}  Prompt sent", tag, CYAN)

    # ── Waiting / streaming phase ──
    t0 = time.time()
    chunks_received = 0
    bytes_received = 0
    text_parts = []
    usage = None
    saved_path = None
    phase = "thinking"
    last_spinner_time = t0

    log(f"{MAGENTA}⏳ Waiting for model{RESET}  {slide_title}", tag, MAGENTA)

    for chunk in client.models.generate_content_stream(
        model=IMAGE_MODEL,
        contents=contents,
        config=config,
    ):
        chunks_received += 1
        elapsed = time.time() - t0
        now = time.time()

        if getattr(chunk, "usage_metadata", None):
            usage = chunk.usage_metadata

        if chunk.parts is None:
            # Model is thinking — show spinner
            if now - last_spinner_time >= 3.0:
                spinner = SPINNER_FRAMES[chunks_received % len(SPINNER_FRAMES)]
                log(
                    f"  {spinner} {DIM}Model is {phase}... "
                    f"({format_duration(elapsed)}, {chunks_received} chunks){RESET}",
                    tag, MAGENTA,
                )
                last_spinner_time = now
            continue

        for part in chunk.parts:
            if part.inline_data and part.inline_data.data:
                phase = "rendering"
                data = part.inline_data.data
                bytes_received += len(data)

                # ── Download complete ──
                dl_bar = progress_bar(1.0, color=GREEN)
                log(
                    f"  {dl_bar}  {GREEN}⬇ Image received{RESET}  "
                    f"{format_size(len(data))}  "
                    f"{DIM}({part.inline_data.mime_type}){RESET}",
                    tag, GREEN,
                )

                ext = mimetypes.guess_extension(part.inline_data.mime_type) or ".png"
                saved_path = next_available_path(output_dir, f"slide_{slide_num:02d}", ext)

                # ── Save phase ──
                save_binary_file(saved_path, data)
                log(
                    f"  {GREEN}💾 Saved{RESET}  {saved_path}  ({format_size(len(data))})",
                    tag, GREEN,
                )

            elif part.text:
                if phase == "thinking":
                    phase = "generating"
                text_parts.append(part.text)
                if now - last_spinner_time >= 3.0:
                    spinner = SPINNER_FRAMES[chunks_received % len(SPINNER_FRAMES)]
                    log(
                        f"  {spinner} {DIM}Model is {phase}... "
                        f"(+{len(part.text)} chars text, {format_duration(elapsed)}){RESET}",
                        tag, MAGENTA,
                    )
                    last_spinner_time = now

    elapsed = time.time() - t0
    slide_cost, slide_breakdown = compute_cost(IMAGE_MODEL, usage)

    # ── Summary ──
    if slide_breakdown:
        log(
            f"  {DIM}Tokens: in:{slide_breakdown['input_tokens']:,}  "
            f"out:{slide_breakdown['output_tokens']:,}  "
            f"think:{slide_breakdown['thinking_tokens']:,}  "
            f"img:{slide_breakdown['image_output_tokens']:,}  "
            f"│  Cost: ${slide_cost:.4f}  │  Time: {format_duration(elapsed)}{RESET}",
            tag, WHITE,
        )

    if saved_path:
        return saved_path, slide_cost

    if text_parts:
        full_text = "".join(text_parts)
        log(f"{RED}✗ Model returned text only:{RESET} {full_text[:300]}...", tag, RED)

    raise RuntimeError(f"No image generated for slide {slide_num}: {slide_title}")


# ── PowerPoint ───────────────────────────────────────────────────────────────

def build_pptx(slides, output_dir, pptx_path):
    """Assemble generated slide images into a PowerPoint file with speaker notes."""
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    blank_layout = prs.slide_layouts[6]

    added = 0
    for slide_data in sorted(slides, key=lambda s: s["slide_number"]):
        slide_num = slide_data["slide_number"]
        image_path = None
        for ext in (".png", ".jpg", ".jpeg", ".webp"):
            candidate = os.path.join(output_dir, f"slide_{slide_num:02d}{ext}")
            if os.path.exists(candidate):
                image_path = candidate
                break

        if not image_path:
            log(f"{YELLOW}⚠ No image for slide {slide_num}, skipping{RESET}", "PPTX", YELLOW)
            continue

        pptx_slide = prs.slides.add_slide(blank_layout)
        pptx_slide.shapes.add_picture(
            image_path, Emu(0), Emu(0),
            width=prs.slide_width, height=prs.slide_height,
        )

        notes = slide_data.get("speaker_notes")
        if notes:
            notes_slide = pptx_slide.notes_slide
            notes_slide.notes_text_frame.text = notes

        added += 1
        log(f"  Added slide {slide_num}  {DIM}({format_size(os.path.getsize(image_path))}){RESET}", "PPTX", WHITE)

    prs.save(pptx_path)
    log(f"{GREEN}✓ PowerPoint saved{RESET}  {pptx_path}  ({added} slides)", "PPTX", GREEN)
    return pptx_path


# ── OCR ──────────────────────────────────────────────────────────────────────

def ocr_slide_image(image_path):
    """Use macOS Vision framework to OCR text from a slide image."""
    import Quartz
    import Vision

    url = Quartz.CFURLCreateWithFileSystemPath(
        None, image_path, Quartz.kCFURLPOSIXPathStyle, False
    )
    source = Quartz.CGImageSourceCreateWithURL(url, None)
    if source is None:
        return []
    image = Quartz.CGImageSourceCreateImageAtIndex(source, 0, None)
    if image is None:
        return []

    request = Vision.VNRecognizeTextRequest.alloc().init()
    request.setRecognitionLevel_(Vision.VNRequestTextRecognitionLevelAccurate)

    handler = Vision.VNImageRequestHandler.alloc().initWithCGImage_options_(image, None)
    success, error = handler.performRequests_error_([request], None)
    if not success:
        return []

    results = request.results()
    lines = []
    for observation in results:
        candidates = observation.topCandidates_(1)
        if candidates:
            lines.append(candidates[0].string())
    return lines


def build_ocr_markdown(slides, output_dir, markdown_path):
    """OCR each slide image and write a formatted markdown file."""
    md_lines = ["# Slide Deck — OCR Transcript\n"]
    ocr_count = 0

    for slide_data in sorted(slides, key=lambda s: s["slide_number"]):
        slide_num = slide_data["slide_number"]
        slide_title = slide_data.get("title", f"Slide {slide_num}")

        image_path = None
        for ext in (".png", ".jpg", ".jpeg", ".webp"):
            candidate = os.path.join(output_dir, f"slide_{slide_num:02d}{ext}")
            if os.path.exists(candidate):
                image_path = candidate
                break

        if not image_path:
            log(f"  {YELLOW}⚠ No image for slide {slide_num}, skipping{RESET}", "OCR", YELLOW)
            continue

        t0 = time.time()
        text_lines = ocr_slide_image(image_path)
        elapsed = time.time() - t0
        word_count = sum(len(line.split()) for line in text_lines)
        log(
            f"  Slide {slide_num}: {len(text_lines)} lines, {word_count} words  "
            f"{DIM}({format_duration(elapsed)}){RESET}",
            "OCR", WHITE,
        )

        md_lines.append(f"---\n\n## Slide {slide_num}: {slide_title}\n")
        if text_lines:
            for line in text_lines:
                md_lines.append(f"{line}\n")
        else:
            md_lines.append("*(no text detected)*\n")
        md_lines.append("")
        ocr_count += 1

    content = "\n".join(md_lines)
    with open(markdown_path, "w") as f:
        f.write(content)
    log(f"{GREEN}✓ OCR markdown saved{RESET}  {markdown_path}  ({ocr_count} slides)", "OCR", GREEN)
    return markdown_path


# ── PDF ──────────────────────────────────────────────────────────────────────

def build_pdf(slides, output_dir, pdf_path):
    """Create a PDF from generated slide images using Pillow."""
    from PIL import Image

    images = []
    for slide_data in sorted(slides, key=lambda s: s["slide_number"]):
        slide_num = slide_data["slide_number"]
        image_path = None
        for ext in (".png", ".jpg", ".jpeg", ".webp"):
            candidate = os.path.join(output_dir, f"slide_{slide_num:02d}{ext}")
            if os.path.exists(candidate):
                image_path = candidate
                break
        if not image_path:
            log(f"  {YELLOW}⚠ No image for slide {slide_num}, skipping{RESET}", "PDF", YELLOW)
            continue
        img = Image.open(image_path).convert("RGB")
        images.append(img)

    if not images:
        log(f"{RED}✗ No images found for PDF{RESET}", "PDF", RED)
        return None

    images[0].save(pdf_path, save_all=True, append_images=images[1:], resolution=150)
    size = format_size(os.path.getsize(pdf_path))
    log(f"{GREEN}✓ PDF saved{RESET}  {pdf_path}  ({len(images)} slides, {size})", "PDF", GREEN)
    return pdf_path


def compress_pdf(input_path, output_path):
    """Compress a PDF using Ghostscript."""
    try:
        subprocess.run(
            [
                "gs", "-sDEVICE=pdfwrite", "-dCompatibilityLevel=1.4",
                "-dPDFSETTINGS=/ebook", "-dNOPAUSE", "-dQUIET", "-dBATCH",
                f"-sOutputFile={output_path}", input_path,
            ],
            check=True, capture_output=True,
        )
        original = os.path.getsize(input_path)
        compressed = os.path.getsize(output_path)
        ratio = (1 - compressed / original) * 100 if original > 0 else 0
        log(
            f"{GREEN}✓ Compressed PDF saved{RESET}  {output_path}  "
            f"({format_size(compressed)}, {ratio:.0f}% smaller)",
            "PDF", GREEN,
        )
        return output_path
    except FileNotFoundError:
        log(f"{YELLOW}⚠ Ghostscript (gs) not found — skipping PDF compression{RESET}", "PDF", YELLOW)
        return None
    except subprocess.CalledProcessError as e:
        log(f"{RED}✗ PDF compression failed: {e}{RESET}", "PDF", RED)
        return None


# ── Parallel generation ──────────────────────────────────────────────────────

def generate_all_slides(client, slides, resolution, output_dir, max_parallel, permutations=1):
    """Generate all slide images in parallel. Returns (paths, total_cost, per_slide_costs)."""
    os.makedirs(output_dir, exist_ok=True)

    tasks = [slide for slide in slides for _ in range(permutations)]
    total = len(tasks)
    completed = 0
    failed = 0
    results = []
    total_cost = 0.0
    per_slide_costs = {}
    failed_slides = []

    log(
        f"Generating {BOLD}{total}{RESET} images  "
        f"({len(slides)} slides × {permutations} variation(s))  "
        f"workers={max_parallel}",
        "GEN", CYAN,
    )
    log(f"Output: {output_dir}", "GEN", CYAN)
    log("", "GEN", WHITE)

    t0 = time.time()

    with concurrent.futures.ThreadPoolExecutor(max_workers=max_parallel) as executor:
        futures = {}
        for slide in tasks:
            f = executor.submit(generate_slide_image, client, slide, resolution, output_dir)
            futures[f] = slide

        for future in concurrent.futures.as_completed(futures):
            slide = futures[future]
            try:
                path, slide_cost = future.result()
                completed += 1
                results.append(path)
                total_cost += slide_cost
                per_slide_costs.setdefault(slide["slide_number"], []).append(slide_cost)
                elapsed = time.time() - t0

                overall_bar = progress_bar((completed + failed) / total, color=GREEN)
                log(
                    f"{overall_bar}  {GREEN}✓{RESET} {os.path.basename(path)}  "
                    f"{DIM}${slide_cost:.4f}  {format_duration(elapsed)} elapsed{RESET}  "
                    f"[{completed}/{total} done]",
                    "GEN", GREEN,
                )

            except Exception as e:
                failed += 1
                elapsed = time.time() - t0
                failed_slides.append(slide)

                overall_bar = progress_bar((completed + failed) / total, color=RED)
                log(
                    f"{overall_bar}  {RED}✗ Slide {slide['slide_number']:02d}{RESET}  "
                    f"{slide['title']}  "
                    f"{DIM}{format_duration(elapsed)} elapsed{RESET}  "
                    f"[{completed}/{total} done, {failed} failed]",
                    "GEN", RED,
                )
                log(f"  {RED}Error: {e}{RESET}", "GEN", RED)

    elapsed = time.time() - t0
    log("", "GEN", WHITE)

    if failed == 0:
        log(
            f"{GREEN}✓ All {completed} images generated successfully{RESET}  "
            f"in {format_duration(elapsed)}  cost=${total_cost:.4f}",
            "GEN", GREEN,
        )
    else:
        log(
            f"{YELLOW}⚠ {completed} succeeded, {failed} failed{RESET}  "
            f"in {format_duration(elapsed)}  cost=${total_cost:.4f}",
            "GEN", YELLOW,
        )
        for s in failed_slides:
            log(f"  {RED}✗ Slide {s['slide_number']}: {s['title']}{RESET}", "GEN", RED)

    return results, total_cost, per_slide_costs


# ── Main ─────────────────────────────────────────────────────────────────────

def main():
    args = parse_args()

    selected_slides = None
    if args.slides:
        selected_slides = parse_slide_numbers(args.slides)

    # ── Startup banner ──
    log_header("Pitch Deck Maker")

    log(f"  {BOLD}Deck file{RESET}     {args.deck_file}", "MAIN", WHITE)
    log(f"  {BOLD}Output dir{RESET}    {args.output}", "MAIN", WHITE)
    log(f"  {BOLD}Resolution{RESET}    {args.resolution}", "MAIN", WHITE)
    log(f"  {BOLD}Workers{RESET}       {args.max_parallel}", "MAIN", WHITE)
    log(f"  {BOLD}Permutations{RESET}  {args.permutations}", "MAIN", WHITE)
    log(f"  {BOLD}Retries{RESET}       {MAX_RETRIES} (backoff: {RETRY_BACKOFF}s)", "MAIN", WHITE)

    if args.onesheet:
        slides_label = "onesheet"
    elif args.all:
        slides_label = "all"
    elif any(isinstance(x, tuple) for x in selected_slides):
        start = next(x[1] for x in selected_slides if isinstance(x, tuple))
        slides_label = f"{start}+"
    else:
        slides_label = str(sorted(selected_slides))
    log(f"  {BOLD}Slides{RESET}        {slides_label}", "MAIN", WHITE)

    if not args.onesheet:
        log(f"  {BOLD}Parse model{RESET}   {PARTITION_MODEL}", "MAIN", WHITE)
    log(f"  {BOLD}Image model{RESET}   {IMAGE_MODEL}", "MAIN", WHITE)

    api_key = os.environ.get("GEMINI_API_KEY")
    if not api_key:
        log(f"{RED}✗ GEMINI_API_KEY environment variable is not set!{RESET}", "ERROR", RED)
        sys.exit(1)
    log(f"  {GREEN}✓ API key loaded{RESET}  ({len(api_key)} chars)", "MAIN", GREEN)

    client = genai.Client(api_key=api_key)
    log(f"  {GREEN}✓ Gemini client initialized{RESET}", "MAIN", GREEN)

    with open(args.deck_file) as f:
        markdown_content = f.read()
    log(
        f"  {GREEN}✓ Deck loaded{RESET}  "
        f"{len(markdown_content):,} chars, {markdown_content.count(chr(10)):,} lines",
        "MAIN", GREEN,
    )

    # ── Phase 1: Parse ──
    if args.onesheet:
        deck_name = os.path.splitext(os.path.basename(args.deck_file))[0]
        slides = [{"slide_number": 1, "title": deck_name, "prompt": markdown_content.strip()}]
        parse_cost = 0.0
        log(f"\n  {CYAN}Onesheet mode — using file as image prompt directly{RESET}", "MAIN", CYAN)
    else:
        log_header("Phase 1: Parsing Deck")
        all_slides, parse_cost = partition_deck(client, markdown_content)
        log(f"\n  Found {BOLD}{len(all_slides)}{RESET} slides with image prompts", "MAIN", WHITE)

        if selected_slides:
            max_slide = max(s["slide_number"] for s in all_slides)
            selected_slides = resolve_slide_range(selected_slides, max_slide)

        if selected_slides:
            slides = [s for s in all_slides if s["slide_number"] in selected_slides]
            missing = selected_slides - {s["slide_number"] for s in all_slides}
            if missing:
                log(f"  {YELLOW}⚠ Slide numbers not found: {sorted(missing)}{RESET}", "MAIN", YELLOW)
            log(f"  Selected {BOLD}{len(slides)}{RESET} of {len(all_slides)} slides", "MAIN", WHITE)
        else:
            slides = all_slides

    # ── Create timestamped batch folder ──
    deck_name = os.path.splitext(os.path.basename(args.deck_file))[0]
    timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
    batch_dir = os.path.join(args.output, f"{deck_name}_{timestamp}")
    os.makedirs(batch_dir, exist_ok=True)
    log(f"  {GREEN}✓ Batch folder{RESET}  {batch_dir}", "MAIN", GREEN)

    # ── Phase 2: Generate ──
    log_header("Phase 2: Generating Images")
    results, gen_cost, per_slide_costs = generate_all_slides(
        client, slides, args.resolution, batch_dir, args.max_parallel, args.permutations
    )

    # ── Phase 3: PowerPoint ──
    if not args.onesheet:
        log_header("Phase 3: Assembling PowerPoint")
        pptx_path = os.path.join(batch_dir, f"{deck_name}.pptx")
        build_pptx(slides, batch_dir, pptx_path)

    # ── Phase 4: PDF ──
    if not args.onesheet:
        log_header("Phase 4: Creating PDF")
        pdf_path = os.path.join(batch_dir, f"{deck_name}.pdf")
        build_pdf(slides, batch_dir, pdf_path)

        compressed_pdf_path = os.path.join(batch_dir, f"{deck_name}_compressed.pdf")
        compress_pdf(pdf_path, compressed_pdf_path)

    # ── Phase 5: OCR ──
    log_header("Phase 5: OCR Slide Images")
    ocr_md_path = os.path.join(batch_dir, f"{deck_name}_ocr.md")
    build_ocr_markdown(slides, batch_dir, ocr_md_path)

    # ── Cost summary ──
    total_cost = parse_cost + gen_cost
    log_header("Cost Summary")

    if not args.onesheet:
        log(f"  Parsing              {BOLD}${parse_cost:.4f}{RESET}", "COST", WHITE)
    log(f"  Image generation     {BOLD}${gen_cost:.4f}{RESET}", "COST", WHITE)

    for slide_num in sorted(per_slide_costs):
        costs = per_slide_costs[slide_num]
        title = next(
            (s["title"] for s in slides if s["slide_number"] == slide_num), ""
        )
        for i, cost in enumerate(costs):
            suffix = f" (variation {i + 1})" if len(costs) > 1 else ""
            log(f"    {DIM}Slide {slide_num:2d}{suffix}: ${cost:.4f}  {title}{RESET}", "COST", WHITE)

    log(f"  {'─' * 36}", "COST", WHITE)
    log(f"  {BOLD}TOTAL                ${total_cost:.4f}{RESET}", "COST", GREEN)

    # ── Done ──
    log_header("Done")
    log(f"  {GREEN}✓{RESET} {BOLD}{len(results)}{RESET} images saved to {batch_dir}/", "DONE", GREEN)
    if not args.onesheet:
        log(f"  {GREEN}✓{RESET} PowerPoint: {pptx_path}", "DONE", GREEN)
        log(f"  {GREEN}✓{RESET} PDF: {pdf_path}", "DONE", GREEN)
        log(f"  {GREEN}✓{RESET} Compressed PDF: {compressed_pdf_path}", "DONE", GREEN)
    log(f"  {GREEN}✓{RESET} OCR markdown: {ocr_md_path}", "DONE", GREEN)


if __name__ == "__main__":
    main()
